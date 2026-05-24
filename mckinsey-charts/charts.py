"""McKinsey-style native python-pptx charts.

Three workhorse types:
  - add_bar_callout:           single-series column with one bar highlighted + a callout
  - add_stacked_bar_over_time: stacked column over time with one series highlighted
  - add_waterfall:             start -> +adds -> -subtracts -> total bridge
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Palette + type system
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x05, 0x1C, 0x2C)
HIGHLIGHT = RGBColor(0x00, 0x6E, 0xB6)   # the one bar / one series that carries the story
GREY = RGBColor(0xBF, 0xBF, 0xBF)        # context bars / series
DARK_GREY = RGBColor(0x55, 0x55, 0x55)   # start / total bars in waterfall
RED = RGBColor(0xC0, 0x39, 0x2B)         # negative deltas in waterfall
SOURCE_GREY = RGBColor(0x80, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Inter"  # falls back to Calibri on systems without Inter

DEFAULT_CHART_BOX = dict(
    left=Inches(0.7),
    top=Inches(1.4),
    width=Inches(11.9),
    height=Inches(5.0),
)


# ---------------------------------------------------------------------------
# Deck + slide helpers
# ---------------------------------------------------------------------------
def new_deck():
    """Return a 16:9 Presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _add_title(slide, title_text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = NAVY


def _add_source(slide, source_text):
    if not source_text:
        return
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Source: {source_text}"
    run.font.name = FONT_NAME
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = SOURCE_GREY


def _add_y_axis_label(slide, label, chart_box):
    if not label:
        return
    tb = slide.shapes.add_textbox(
        Inches(0.5),
        chart_box["top"] - Inches(0.05),
        Inches(3),
        Inches(0.3),
    )
    p = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.name = FONT_NAME
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = DARK_GREY


# ---------------------------------------------------------------------------
# Generic chart styling
# ---------------------------------------------------------------------------
def _style_chart(chart, show_legend=False):
    """Apply McKinsey-ish defaults: no gridlines, no chart border, optional legend."""
    chart.has_title = False

    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = FONT_NAME
        chart.legend.font.size = Pt(10)
        chart.legend.font.color.rgb = DARK_GREY
    else:
        chart.has_legend = False

    # Hide major gridlines
    try:
        gridlines = chart.value_axis.major_gridlines
        gridlines.format.line.fill.background()
    except Exception:
        pass

    # Value axis: light grey ticks, no axis line
    va = chart.value_axis
    va.tick_labels.font.name = FONT_NAME
    va.tick_labels.font.size = Pt(10)
    va.tick_labels.font.color.rgb = DARK_GREY
    try:
        va.format.line.fill.background()
    except Exception:
        pass

    # Category axis: keep visible, light styling
    ca = chart.category_axis
    ca.tick_labels.font.name = FONT_NAME
    ca.tick_labels.font.size = Pt(10)
    ca.tick_labels.font.color.rgb = DARK_GREY


def _set_point_color(point, rgb):
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    # also kill the outline for a flat look
    line = point.format.line
    line.fill.background()


def _set_series_color(series, rgb):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    line = series.format.line
    line.fill.background()


def _hide_series(series):
    """Make a series invisible (used as the floor in the waterfall trick)."""
    fill = series.format.fill
    fill.background()
    line = series.format.line
    line.fill.background()


# ---------------------------------------------------------------------------
# Chart 1: Bar + callout
# ---------------------------------------------------------------------------
def add_bar_callout(
    slide,
    title,
    categories,
    values,
    highlight_index=None,
    callout=None,
    y_label=None,
    source=None,
    chart_box=None,
):
    """Single-series column chart with one bar highlighted and an optional callout textbox.

    Use for TAM-style stories where one number is the headline.
    """
    if highlight_index is None:
        highlight_index = len(values) - 1
    box = chart_box or DEFAULT_CHART_BOX

    _add_title(slide, title)

    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    chart_data.add_series("Value", list(values))

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        box["left"], box["top"], box["width"], box["height"],
        chart_data,
    )
    chart = graphic_frame.chart
    _style_chart(chart, show_legend=False)

    plot = chart.plots[0]
    plot.gap_width = 80  # tighter bars look more confident

    series = plot.series[0]
    _set_series_color(series, GREY)

    # Highlight the chosen bar
    for i in range(len(values)):
        try:
            pt = series.points[i]
            _set_point_color(pt, HIGHLIGHT if i == highlight_index else GREY)
        except IndexError:
            pass

    # Data labels above bars
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.name = FONT_NAME
    dl.font.size = Pt(10)
    dl.font.color.rgb = DARK_GREY
    dl.font.bold = False

    _add_y_axis_label(slide, y_label, box)
    _add_source(slide, source)

    if callout:
        _add_callout(slide, callout, anchor_box=box, anchor_fraction=highlight_index / max(len(values) - 1, 1))

    return chart


def _add_callout(slide, text, anchor_box, anchor_fraction):
    """Drop a bold callout textbox above the highlighted bar."""
    chart_left = anchor_box["left"]
    chart_width = anchor_box["width"]
    inner_left = chart_left + Inches(0.7)             # crude offset past Y axis
    inner_width = chart_width - Inches(1.0)
    callout_width = Inches(2.6)
    # center horizontally over the anchor bar
    bar_center = inner_left + Emu(int(inner_width * anchor_fraction))
    cleft = bar_center - Emu(int(callout_width / 2))
    ctop = anchor_box["top"] + Inches(0.1)
    tb = slide.shapes.add_textbox(cleft, ctop, callout_width, Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = FONT_NAME
        run.font.size = Pt(22) if i == 0 else Pt(11)
        run.font.bold = (i == 0)
        run.font.color.rgb = HIGHLIGHT if i == 0 else DARK_GREY


# ---------------------------------------------------------------------------
# Chart 2: Stacked bar over time
# ---------------------------------------------------------------------------
def add_stacked_bar_over_time(
    slide,
    title,
    categories,
    series,
    highlight_series=None,
    y_label=None,
    source=None,
    chart_box=None,
):
    """Stacked column over time. `series` is a list of (name, values) tuples.

    If highlight_series is set, that series uses HIGHLIGHT color and others use grey shades.
    """
    box = chart_box or DEFAULT_CHART_BOX
    _add_title(slide, title)

    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for name, values in series:
        chart_data.add_series(name, list(values))

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        box["left"], box["top"], box["width"], box["height"],
        chart_data,
    )
    chart = graphic_frame.chart
    _style_chart(chart, show_legend=True)

    plot = chart.plots[0]
    plot.gap_width = 80

    # Grey ramp for non-highlight series
    grey_ramp = [
        RGBColor(0xD9, 0xD9, 0xD9),
        RGBColor(0xBF, 0xBF, 0xBF),
        RGBColor(0xA6, 0xA6, 0xA6),
        RGBColor(0x80, 0x80, 0x80),
    ]
    grey_idx = 0
    for s in plot.series:
        if highlight_series and s.name == highlight_series:
            _set_series_color(s, HIGHLIGHT)
        else:
            _set_series_color(s, grey_ramp[grey_idx % len(grey_ramp)])
            grey_idx += 1

    _add_y_axis_label(slide, y_label, box)
    _add_source(slide, source)
    return chart


# ---------------------------------------------------------------------------
# Chart 3: Waterfall (stacked-bar trick)
# ---------------------------------------------------------------------------
def add_waterfall(
    slide,
    title,
    labels,
    values,
    kinds,
    y_label=None,
    source=None,
    chart_box=None,
):
    """Waterfall built from stacked columns.

    kinds[i] ∈ {"start", "pos", "neg", "total"}
    values[i] is the *signed* delta for "pos"/"neg" and the *absolute* value for "start"/"total".

    Example:
        labels = ["FY24", "New logos", "Expansion", "Churn", "Price", "FY25"]
        values = [80,      24,           12,          -8,      4,       112]
        kinds  = ["start", "pos",        "pos",       "neg",   "pos",   "total"]
    """
    if not (len(labels) == len(values) == len(kinds)):
        raise ValueError("labels, values, kinds must be same length")

    box = chart_box or DEFAULT_CHART_BOX
    _add_title(slide, title)

    # Compute base + visible heights per bar
    bases = []
    visibles = {"start": [], "pos": [], "neg": [], "total": []}
    running = 0.0
    for label, v, k in zip(labels, values, kinds):
        # default: nothing in this row for series we're not in
        for key in visibles:
            visibles[key].append(0)

        if k == "start":
            bases.append(0)
            visibles["start"][-1] = v
            running = v
        elif k == "pos":
            bases.append(running)
            visibles["pos"][-1] = v
            running += v
        elif k == "neg":
            # v is negative; bar drops from running to running+v
            running += v
            bases.append(running)
            visibles["neg"][-1] = -v  # absolute height
        elif k == "total":
            bases.append(0)
            visibles["total"][-1] = v
            # leave running as-is (sanity check could compare to v)
        else:
            raise ValueError(f"Unknown kind: {k!r}")

    chart_data = CategoryChartData()
    chart_data.categories = list(labels)
    chart_data.add_series("__base__", bases)
    chart_data.add_series("Start/End", visibles["start"])  # rendered same color as total
    chart_data.add_series("Increase", visibles["pos"])
    chart_data.add_series("Decrease", visibles["neg"])
    chart_data.add_series("Total", visibles["total"])

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        box["left"], box["top"], box["width"], box["height"],
        chart_data,
    )
    chart = graphic_frame.chart
    _style_chart(chart, show_legend=True)

    plot = chart.plots[0]
    plot.gap_width = 60

    series_map = {s.name: s for s in plot.series}
    _hide_series(series_map["__base__"])
    _set_series_color(series_map["Start/End"], DARK_GREY)
    _set_series_color(series_map["Increase"], HIGHLIGHT)
    _set_series_color(series_map["Decrease"], RED)
    _set_series_color(series_map["Total"], DARK_GREY)

    # Hide __base__ from the legend by deleting its legend entry isn't supported via API;
    # cleanest workaround: rename it to a blank-ish label so it's visually inoffensive,
    # OR skip the legend and rely on color intuition. We skip cleanly here:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = FONT_NAME
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = DARK_GREY
    _delete_legend_entry_for_series(chart, "__base__")

    _add_y_axis_label(slide, y_label, box)
    _add_source(slide, source)
    return chart


def _delete_legend_entry_for_series(chart, series_name):
    """Remove a single legend entry by name (raw XML — python-pptx doesn't expose this)."""
    from pptx.oxml.ns import qn

    plotArea = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    series_idx = None
    for i, ser in enumerate(plotArea.iter(qn("c:ser"))):
        tx = ser.find(qn("c:tx"))
        if tx is None:
            continue
        v = tx.find(".//" + qn("c:v"))
        if v is not None and v.text == series_name:
            idx_el = ser.find(qn("c:idx"))
            if idx_el is not None:
                series_idx = int(idx_el.get("val"))
            break
    if series_idx is None:
        return

    legend = chart._chartSpace.find(qn("c:chart")).find(qn("c:legend"))
    if legend is None:
        return

    # Add a legendEntry that deletes the matching index
    legendEntry = legend.makeelement(qn("c:legendEntry"), {})
    idx = legend.makeelement(qn("c:idx"), {"val": str(series_idx)})
    delete = legend.makeelement(qn("c:delete"), {"val": "1"})
    legendEntry.append(idx)
    legendEntry.append(delete)
    # insert before <c:overlay> if present, else just append
    overlay = legend.find(qn("c:overlay"))
    if overlay is not None:
        overlay.addprevious(legendEntry)
    else:
        legend.append(legendEntry)
